"""Parse PDF / DOCX / PPTX / plain text into (page_ref, text) segments and
detect chapter/section headings where possible."""
import io
import re


def _detect_chapters(full_text: str):
    chapters = []
    for line in full_text.splitlines():
        s = line.strip()
        if not s or len(s) > 90:
            continue
        if re.match(r"^(chapter|unit|section|lesson|part)\s+[\dIVX]+", s, re.I) or \
           re.match(r"^\d+(\.\d+)*\s+[A-Z].{2,}", s):
            chapters.append(s)
    return chapters[:40]


def parse_pdf(data: bytes):
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages):
        txt = page.extract_text() or ""
        if txt.strip():
            pages.append((f"page {i + 1}", txt))
    return pages


def parse_docx(data: bytes):
    from docx import Document
    doc = Document(io.BytesIO(data))
    pages, buf, sec = [], [], 1
    for p in doc.paragraphs:
        t = p.text.strip()
        if not t:
            continue
        buf.append(t)
        if len("\n".join(buf)) > 1200:
            pages.append((f"section {sec}", "\n".join(buf)))
            buf, sec = [], sec + 1
    if buf:
        pages.append((f"section {sec}", "\n".join(buf)))
    return pages


def parse_pptx(data: bytes):
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if parts:
            pages.append((f"slide {i + 1}", "\n".join(parts)))
    return pages


def parse_text(data: bytes):
    text = data.decode("utf-8", errors="ignore")
    pages, chunks = [], text.split("\n\n")
    buf, sec = [], 1
    for c in chunks:
        buf.append(c)
        if len("\n\n".join(buf)) > 1200:
            pages.append((f"section {sec}", "\n\n".join(buf)))
            buf, sec = [], sec + 1
    if buf:
        pages.append((f"section {sec}", "\n\n".join(buf)))
    return pages


def parse_file(filename: str, data: bytes):
    name = (filename or "").lower()
    if name.endswith(".pdf"):
        source_type, pages = "pdf", parse_pdf(data)
    elif name.endswith(".docx"):
        source_type, pages = "docx", parse_docx(data)
    elif name.endswith(".pptx"):
        source_type, pages = "pptx", parse_pptx(data)
    else:
        source_type, pages = "text", parse_text(data)
    full_text = "\n".join(t for _, t in pages)
    return {
        "source_type": source_type,
        "pages": pages,
        "raw_text": full_text[:200000],
        "chapters": _detect_chapters(full_text),
    }


def chunk_pages(pages, size: int = 900, overlap: int = 150):
    chunks = []
    for page_ref, text in pages:
        text = re.sub(r"\s+", " ", text).strip()
        i = 0
        while i < len(text):
            piece = text[i:i + size]
            if piece.strip():
                chunks.append({"page_ref": page_ref, "chunk_text": piece})
            i += size - overlap
    return chunks
